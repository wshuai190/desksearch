"""Document parsers for various file formats.

Registry-based parser system. Each parser is a callable that takes a file path
and returns extracted plain text. New parsers can be added via register_parser().
"""
import logging
import time
from pathlib import Path
from typing import Callable, Optional

logger = logging.getLogger(__name__)

# Parser type: takes a Path, returns extracted text
ParserFunc = Callable[[Path], str]

# Registry mapping extensions to parser functions
_PARSERS: dict[str, ParserFunc] = {}

# Code file extensions for language detection
CODE_EXTENSIONS: dict[str, str] = {
    ".py": "python",
    ".js": "javascript",
    ".ts": "typescript",
    ".java": "java",
    ".c": "c",
    ".cpp": "cpp",
    ".h": "c-header",
    ".go": "go",
    ".rs": "rust",
    ".sh": "bash",
    ".bash": "bash",
    ".zsh": "zsh",
    ".sql": "sql",
    ".r": "r",
    ".R": "r",
    ".rb": "ruby",
    ".swift": "swift",
    ".kt": "kotlin",
    ".scala": "scala",
    ".lua": "lua",
    ".pl": "perl",
    ".php": "php",
}


def register_parser(extensions: list[str], parser: ParserFunc) -> None:
    """Register a parser function for one or more file extensions."""
    for ext in extensions:
        _PARSERS[ext.lower()] = parser


def get_parser(extension: str) -> Optional[ParserFunc]:
    """Get the parser for a given file extension."""
    return _PARSERS.get(extension.lower())


def parse_file(path: Path) -> Optional[str]:
    """Parse a file and return its text content.

    Args:
        path: Path to the file to parse.

    Returns:
        Extracted text, or None if the file cannot be parsed.
    """
    ext = path.suffix.lower()
    parser = get_parser(ext)
    if parser is None:
        logger.warning("No parser registered for extension: %s", ext)
        return None
    t0 = time.perf_counter()
    try:
        text = parser(path)
        elapsed_ms = (time.perf_counter() - t0) * 1000
        parser_name = getattr(parser, "__name__", "unknown")
        if text and text.strip():
            logger.debug(
                "[%s] parser=%s elapsed=%.1fms chars=%d",
                path.name, parser_name, elapsed_ms, len(text),
            )
            return text.strip()
        logger.warning("Parser returned empty text for: %s (%.1fms)", path, elapsed_ms)
        return None
    except Exception:
        elapsed_ms = (time.perf_counter() - t0) * 1000
        logger.exception("Failed to parse file: %s (%.1fms)", path, elapsed_ms)
        return None


# --- Built-in parsers ---


def _parse_text(path: Path) -> str:
    """Parse plain text files (txt, md, rst, org, csv, tsv, etc.)."""
    return path.read_text(encoding="utf-8", errors="replace")


def _parse_pdf(path: Path) -> str:
    """Parse PDF files using PyMuPDF (fitz)."""
    import fitz

    pages = []
    with fitz.open(str(path)) as doc:
        for page in doc:
            text = page.get_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def _parse_docx(path: Path) -> str:
    """Parse DOCX files using python-docx."""
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _parse_code(path: Path) -> str:
    """Parse code files with language metadata prefix."""
    ext = path.suffix.lower()
    language = CODE_EXTENSIONS.get(ext, "unknown")
    content = path.read_text(encoding="utf-8", errors="replace")
    return f"[{language}] {path.name}\n\n{content}"


def _parse_html(path: Path) -> str:
    """Parse HTML files, extracting text content."""
    from bs4 import BeautifulSoup

    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_json_yaml(path: Path) -> str:
    """Parse JSON/YAML/TOML as plain text (structure is content)."""
    return path.read_text(encoding="utf-8", errors="replace")


def _parse_ipynb(path: Path) -> str:
    """Parse Jupyter notebooks, extracting markdown and code cells."""
    import json

    data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    parts = []
    for cell in data.get("cells", []):
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", []))
        if cell_type == "markdown":
            parts.append(source)
        elif cell_type == "code":
            parts.append(f"```\n{source}\n```")
    return "\n\n".join(parts)


def _parse_pptx(path: Path) -> str:
    """Parse PowerPoint files using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _parse_xlsx(path: Path) -> str:
    """Parse Excel files using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[{sheet.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _parse_epub(path: Path) -> str:
    """Parse EPUB files by extracting HTML content from the archive."""
    import zipfile
    from bs4 import BeautifulSoup

    parts = []
    with zipfile.ZipFile(str(path), "r") as zf:
        for name in zf.namelist():
            if name.endswith((".html", ".xhtml", ".htm")):
                raw = zf.read(name).decode("utf-8", errors="replace")
                soup = BeautifulSoup(raw, "html.parser")
                for tag in soup(["script", "style"]):
                    tag.decompose()
                text = soup.get_text(separator="\n", strip=True)
                if text.strip():
                    parts.append(text)
    return "\n\n".join(parts)


def _parse_rtf(path: Path) -> str:
    """Parse RTF files by stripping RTF control words."""
    import re
    raw = path.read_text(encoding="utf-8", errors="replace")
    # Strip RTF control words and groups
    text = re.sub(r'\\[a-z]+\d*\s?', ' ', raw)
    text = re.sub(r'[{}]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _parse_archive(path: Path) -> str:
    """Parse archive files (zip/tar/gz) by extracting and parsing contained text files."""
    import zipfile
    import tarfile
    import io

    text_exts = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs", ".sh", ".sql", ".toml", ".rst", ".org", ".tex", ".r", ".log", ".cfg", ".ini", ".conf"}
    parts = []
    max_files = 100  # don't go crazy inside huge archives
    count = 0

    ext = path.suffix.lower()
    try:
        if ext == ".zip":
            with zipfile.ZipFile(str(path), "r") as zf:
                for info in zf.infolist():
                    if count >= max_files:
                        break
                    if info.is_dir():
                        continue
                    inner_ext = Path(info.filename).suffix.lower()
                    if inner_ext in text_exts:
                        try:
                            raw = zf.read(info.filename).decode("utf-8", errors="replace")
                            if raw.strip():
                                parts.append(f"[{info.filename}]\n{raw[:50000]}")
                                count += 1
                        except Exception:
                            continue
        elif ext in (".tar", ".gz", ".tgz", ".bz2", ".xz"):
            mode = "r:*" if ext != ".tar" else "r"
            with tarfile.open(str(path), mode) as tf:
                for member in tf.getmembers():
                    if count >= max_files:
                        break
                    if not member.isfile():
                        continue
                    inner_ext = Path(member.name).suffix.lower()
                    if inner_ext in text_exts:
                        try:
                            f = tf.extractfile(member)
                            if f:
                                try:
                                    raw = f.read().decode("utf-8", errors="replace")
                                finally:
                                    f.close()  # Explicitly release the stream
                                if raw.strip():
                                    parts.append(f"[{member.name}]\n{raw[:50000]}")
                                    count += 1
                        except Exception:
                            continue
    except Exception as e:
        logger.warning("Failed to read archive %s: %s", path, e)

    return "\n\n".join(parts) if parts else ""


def _parse_eml(path: Path) -> str:
    """Parse email files (.eml)."""
    import email
    from email import policy

    raw = path.read_bytes()
    msg = email.message_from_bytes(raw, policy=policy.default)
    parts = []
    subject = msg.get("subject", "")
    sender = msg.get("from", "")
    date = msg.get("date", "")
    if subject:
        parts.append(f"Subject: {subject}")
    if sender:
        parts.append(f"From: {sender}")
    if date:
        parts.append(f"Date: {date}")
    parts.append("")

    body = msg.get_body(preferencelist=("plain", "html"))
    if body:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, "html.parser")
            content = soup.get_text(separator="\n", strip=True)
        parts.append(content)

    return "\n".join(parts)


# --- Register built-in parsers ---

register_parser([".txt", ".md", ".rst", ".org", ".tex"], _parse_text)
register_parser([".csv", ".tsv"], _parse_text)
register_parser([".log", ".cfg", ".ini", ".conf", ".env"], _parse_text)
register_parser([".pdf"], _parse_pdf)
register_parser([".docx", ".doc"], _parse_docx)
register_parser([".pptx", ".ppt"], _parse_pptx)
register_parser([".xlsx", ".xls"], _parse_xlsx)
register_parser([".epub"], _parse_epub)
register_parser([".rtf"], _parse_rtf)
register_parser([".html", ".htm", ".xml"], _parse_html)
register_parser([".json", ".yaml", ".yml", ".toml"], _parse_json_yaml)
register_parser([".ipynb"], _parse_ipynb)
register_parser([".eml", ".msg"], _parse_eml)
register_parser([".zip"], _parse_archive)
register_parser([".tar", ".gz", ".tgz", ".bz2", ".xz"], _parse_archive)
register_parser(list(CODE_EXTENSIONS.keys()), _parse_code)
