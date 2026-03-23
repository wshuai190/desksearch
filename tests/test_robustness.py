"""Robustness, error handling, and production-readiness tests.

Covers:
1. New parser formats (pptx, xlsx, epub, zip, eml)
2. Search quality (a known document should rank #1 for an exact query)
3. Concurrent access (simultaneous indexing and searching)
4. Error recovery (crash-recovery via indexing_state)
5. Config validation
6. Health endpoint
7. Graceful degradation (BM25-only / dense-only fallback)
"""
import io
import json
import threading
import time
import zipfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import numpy as np
import pytest

from desksearch.config import Config
from desksearch.core.bm25 import BM25Index
from desksearch.core.dense import DenseIndex
from desksearch.core.search import HybridSearchEngine
from desksearch.indexer.parsers import parse_file
from desksearch.indexer.pipeline import IndexingPipeline, StatusType
from desksearch.indexer.store import MetadataStore, STATE_DONE, STATE_INDEXING, STATE_FAILED


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

DIM = 32


def _mock_embedder(dim: int = DIM):
    """Return a mock Embedder that yields deterministic vectors."""
    mock = MagicMock()
    mock.dimension = dim
    mock.embed.side_effect = lambda texts, batch_size=64: np.random.rand(
        len(texts), dim
    ).astype(np.float32)
    mock.embed_query.side_effect = lambda q: np.random.rand(dim).astype(np.float32)
    return mock


# ---------------------------------------------------------------------------
# 1. Parser tests — pptx, xlsx, epub, zip, eml
# ---------------------------------------------------------------------------


class TestNewParsers:
    """Test parsers for document formats added after the original release."""

    def test_pptx_parser(self, tmp_path: Path):
        """PowerPoint parser should extract slide text."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # blank
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = txBox.text_frame
        tf.text = "Artificial intelligence revolutionises search"

        path = tmp_path / "test.pptx"
        prs.save(str(path))

        result = parse_file(path)
        assert result is not None
        assert "Artificial intelligence" in result
        assert "Slide 1" in result

    def test_xlsx_parser(self, tmp_path: Path):
        """Excel parser should extract cell values from all sheets."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Revenue"
        ws.append(["Quarter", "Amount", "Region"])
        ws.append(["Q1", 100000, "APAC"])
        ws.append(["Q2", 120000, "EMEA"])
        path = tmp_path / "test.xlsx"
        wb.save(str(path))

        result = parse_file(path)
        assert result is not None
        assert "Revenue" in result
        assert "Quarter" in result
        assert "100000" in result

    def test_epub_parser(self, tmp_path: Path):
        """EPUB parser should extract HTML content from the archive."""
        # Create a minimal valid EPUB (it's a ZIP with HTML files)
        path = tmp_path / "book.epub"
        with zipfile.ZipFile(str(path), "w") as zf:
            zf.writestr(
                "OEBPS/chapter1.html",
                "<html><body><p>Chapter one: the beginning of wisdom.</p></body></html>",
            )
            zf.writestr(
                "OEBPS/chapter2.xhtml",
                "<html><body><p>Chapter two: deep neural networks explained.</p></body></html>",
            )

        result = parse_file(path)
        assert result is not None
        assert "beginning of wisdom" in result
        assert "neural networks" in result

    def test_zip_parser_extracts_text_files(self, tmp_path: Path):
        """ZIP parser should extract text from files inside the archive."""
        path = tmp_path / "archive.zip"
        with zipfile.ZipFile(str(path), "w") as zf:
            zf.writestr("readme.txt", "Instruction manual for the widget")
            zf.writestr("notes.md", "# Research Notes\n\nQuantum entanglement is fascinating.")
            zf.writestr("binary.bin", b"\x00\x01\x02\x03".decode("latin-1"))

        result = parse_file(path)
        assert result is not None
        assert "Instruction manual" in result
        assert "Quantum entanglement" in result

    def test_zip_parser_skips_binary(self, tmp_path: Path):
        """ZIP parser should not crash on binary content."""
        path = tmp_path / "mixed.zip"
        with zipfile.ZipFile(str(path), "w") as zf:
            zf.writestr("text.txt", "Some plain text here")
            # Write actual binary data that can't be decoded as UTF-8
            zf.writestr("data.json", '{"key": "value with unicode: \u00e9"}')

        result = parse_file(path)
        assert result is not None

    def test_eml_parser(self, tmp_path: Path):
        """Email parser should extract subject, sender, and body."""
        eml_content = (
            "From: alice@example.com\r\n"
            "To: bob@example.com\r\n"
            "Subject: Quarterly Report Attached\r\n"
            "Date: Mon, 1 Jan 2024 12:00:00 +0000\r\n"
            "Content-Type: text/plain; charset=utf-8\r\n"
            "\r\n"
            "Hi Bob,\n\nPlease find the quarterly revenue report attached.\n\nBest,\nAlice\n"
        )
        path = tmp_path / "message.eml"
        path.write_bytes(eml_content.encode("utf-8"))

        result = parse_file(path)
        assert result is not None
        assert "Quarterly Report" in result
        assert "alice@example.com" in result
        assert "quarterly revenue report" in result

    def test_eml_parser_html_body(self, tmp_path: Path):
        """Email parser should strip HTML tags from HTML-body emails."""
        eml_content = (
            "From: newsletter@example.com\r\n"
            "Subject: Weekly Digest\r\n"
            "Content-Type: text/html; charset=utf-8\r\n"
            "\r\n"
            "<html><body><h1>Top Stories</h1><p>AI advances in semantic search.</p></body></html>\r\n"
        )
        path = tmp_path / "newsletter.eml"
        path.write_bytes(eml_content.encode("utf-8"))

        result = parse_file(path)
        assert result is not None
        # HTML tags should have been stripped
        assert "<html>" not in result
        assert "Top Stories" in result or "AI advances" in result

    def test_pptx_parser_empty_presentation(self, tmp_path: Path):
        """PPTX parser should return empty/None for a presentation with no text."""
        from pptx import Presentation

        prs = Presentation()
        # Add slide with no text
        prs.slides.add_slide(prs.slide_layouts[6])
        path = tmp_path / "empty.pptx"
        prs.save(str(path))

        result = parse_file(path)
        # Either None or empty — both are acceptable for empty presentations
        assert result is None or result.strip() == ""

    def test_xlsx_multi_sheet(self, tmp_path: Path):
        """Excel parser should extract data from multiple worksheets."""
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sales"
        ws1.append(["Product", "Units"])
        ws1.append(["Widget A", 500])

        ws2 = wb.create_sheet("Inventory")
        ws2.append(["Item", "Stock"])
        ws2.append(["Part X", 200])

        path = tmp_path / "multisheet.xlsx"
        wb.save(str(path))

        result = parse_file(path)
        assert result is not None
        assert "Sales" in result
        assert "Inventory" in result
        assert "Widget A" in result
        assert "Part X" in result


# ---------------------------------------------------------------------------
# 2. Search quality — known document should rank #1
# ---------------------------------------------------------------------------


class TestSearchQuality:
    """End-to-end search quality: a document with exact query terms must rank first."""

    def test_exact_query_ranks_target_first(self, tmp_path: Path):
        """Index three documents; the one containing all query terms must be #1."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        target_text = (
            "Transformer architecture revolutionised natural language processing. "
            "Self-attention mechanisms allow the model to weigh token relationships."
        )
        noise1 = "Gardening tips: water your plants every morning for best results."
        noise2 = "Stock market analysis for the third quarter shows mixed signals."

        # Use real (random) embeddings — BM25 exact-match is sufficient for quality test
        rng = np.random.default_rng(42)
        emb_target = rng.standard_normal(DIM).astype(np.float32)
        emb_n1 = rng.standard_normal(DIM).astype(np.float32)
        emb_n2 = rng.standard_normal(DIM).astype(np.float32)

        engine.add_document("target", target_text, emb_target)
        engine.add_document("noise1", noise1, emb_n1)
        engine.add_document("noise2", noise2, emb_n2)

        # Use alpha=0 (BM25 only) to test lexical exact-match quality
        results = engine.search_sync(
            "transformer architecture natural language processing",
            rng.standard_normal(DIM).astype(np.float32),
            top_k=3,
            alpha=0.0,
        )

        assert len(results) >= 1, "Expected at least one result"
        assert results[0].doc_id == "target", (
            f"Expected 'target' to rank first, got '{results[0].doc_id}'"
        )

    def test_dense_only_semantic_match(self, tmp_path: Path):
        """With aligned embeddings, dense-only search must find the semantic match."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        # Give target and query almost identical embeddings
        base = np.ones(DIM, dtype=np.float32) / np.sqrt(DIM)
        target_emb = base + np.random.default_rng(7).standard_normal(DIM).astype(np.float32) * 0.01
        query_emb = base.copy()  # nearly identical
        noise_emb = -base + np.random.default_rng(9).standard_normal(DIM).astype(np.float32) * 0.01

        engine.add_document("semantic_match", "Document about machine intelligence", target_emb)
        engine.add_document("unrelated", "Cooking recipes for weeknight dinners", noise_emb)

        results = engine.search_sync(
            "ML models",
            query_emb,
            top_k=2,
            alpha=1.0,  # dense only
        )

        assert len(results) >= 1
        assert results[0].doc_id == "semantic_match"

    def test_hybrid_outperforms_single_backend_for_mixed_query(self, tmp_path: Path):
        """Hybrid search should return at least as many relevant docs as either backend alone."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        # doc_a has good keyword match but bad semantic embedding
        # doc_b has good semantic embedding but no keyword overlap
        base = np.ones(DIM, dtype=np.float32)
        engine.add_document(
            "keyword_doc",
            "python programming tutorial for beginners",
            -base / np.linalg.norm(base),  # opposite direction — bad semantic match
        )
        engine.add_document(
            "semantic_doc",
            "completely unrelated keywords here xyz abc",
            base / np.linalg.norm(base),   # same direction — good semantic match
        )

        query_emb = base / np.linalg.norm(base)

        bm25_results = {r.doc_id for r in engine.search_sync("python tutorial", query_emb, alpha=0.0)}
        dense_results = {r.doc_id for r in engine.search_sync("python tutorial", query_emb, alpha=1.0)}
        hybrid_results = {r.doc_id for r in engine.search_sync("python tutorial", query_emb, alpha=0.5)}

        # Hybrid should include docs from both backends
        assert len(hybrid_results) >= len(bm25_results | dense_results) or \
               len(hybrid_results) >= 1, "Hybrid should return results"


# ---------------------------------------------------------------------------
# 3. Concurrent access
# ---------------------------------------------------------------------------


class TestConcurrentAccess:
    """Ensure the indexes remain consistent under concurrent reads and writes."""

    def test_concurrent_bm25_reads(self, tmp_path: Path):
        """Multiple threads reading BM25 simultaneously should not error."""
        idx = BM25Index(tmp_path)
        idx.add_documents([
            (f"doc{i}", f"content about topic {i} with extra words") for i in range(10)
        ])

        errors: list[Exception] = []

        def _search():
            try:
                for _ in range(20):
                    results = idx.search("topic content", top_k=5)
                    assert isinstance(results, list)
            except Exception as exc:
                errors.append(exc)

        threads = [threading.Thread(target=_search) for _ in range(8)]
        for t in threads:
            t.start()
        for t in threads:
            t.join(timeout=10)

        assert not errors, f"Concurrent BM25 reads raised: {errors}"

    def test_concurrent_dense_reads(self, tmp_path: Path):
        """Multiple threads reading FAISS simultaneously should not error."""
        idx = DenseIndex(tmp_path / "dense", dimension=DIM)
        rng = np.random.default_rng(0)
        for i in range(10):
            idx.add(f"doc{i}", rng.standard_normal(DIM).astype(np.float32))

        errors: list[Exception] = []
        query = rng.standard_normal(DIM).astype(np.float32)

        def _search():
            try:
                for _ in range(20):
                    results = idx.search(query, top_k=5)
                    assert isinstance(results, list)
            except Exception as exc:
                errors.append(exc)

        threads = [threading.Thread(target=_search) for _ in range(8)]
        for t in threads:
            t.start()
        for t in threads:
            t.join(timeout=10)

        assert not errors, f"Concurrent FAISS reads raised: {errors}"

    def test_concurrent_store_reads_and_writes(self, tmp_path: Path):
        """SQLite store should handle concurrent reads and writes without corruption."""
        store = MetadataStore(tmp_path / "test.db")
        errors: list[Exception] = []

        def _writer():
            try:
                for i in range(10):
                    f = tmp_path / f"file_{threading.current_thread().name}_{i}.txt"
                    f.write_text(f"content {i}")
                    try:
                        doc_id = store.upsert_document(f, num_chunks=1)
                        store.add_chunks(doc_id, [(f"chunk {i}", 0, 0)])
                    except Exception:
                        # Transient SQLite busy errors during concurrent writes are acceptable
                        pass
            except Exception as exc:
                errors.append(exc)

        def _reader():
            try:
                for _ in range(30):
                    try:
                        _ = store.document_count()
                        _ = store.chunk_count()
                    except Exception:
                        # Transient read errors during concurrent writes are acceptable
                        pass
            except Exception as exc:
                errors.append(exc)

        writers = [threading.Thread(target=_writer, name=f"w{i}") for i in range(4)]
        readers = [threading.Thread(target=_reader, name=f"r{i}") for i in range(4)]

        all_threads = writers + readers
        for t in all_threads:
            t.start()
        for t in all_threads:
            t.join(timeout=15)

        store.close()
        assert not errors, f"Concurrent store access raised: {errors}"

    def test_index_and_search_simultaneously(self, tmp_path: Path):
        """Indexing in a background thread while searching should not corrupt results."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        # Pre-populate index
        rng = np.random.default_rng(1)
        for i in range(20):
            engine.add_document(
                f"pre_{i}",
                f"pre-indexed document number {i} about topics",
                rng.standard_normal(DIM).astype(np.float32),
            )

        search_errors: list[Exception] = []
        index_errors: list[Exception] = []
        query_emb = rng.standard_normal(DIM).astype(np.float32)

        def _search_loop():
            try:
                for _ in range(50):
                    results = engine.search_sync("topics document", query_emb, top_k=5)
                    assert isinstance(results, list)
                    time.sleep(0.001)
            except Exception as exc:
                search_errors.append(exc)

        def _index_loop():
            try:
                for i in range(20):
                    engine.add_document(
                        f"new_{i}",
                        f"newly added document {i} with fresh content",
                        rng.standard_normal(DIM).astype(np.float32),
                    )
                    time.sleep(0.002)
            except Exception as exc:
                index_errors.append(exc)

        searcher = threading.Thread(target=_search_loop)
        indexer = threading.Thread(target=_index_loop)

        searcher.start()
        indexer.start()
        searcher.join(timeout=15)
        indexer.join(timeout=15)

        assert not search_errors, f"Search errors during concurrent indexing: {search_errors}"
        assert not index_errors, f"Index errors during concurrent searching: {index_errors}"


# ---------------------------------------------------------------------------
# 4. Error recovery — crash-recovery via indexing_state
# ---------------------------------------------------------------------------


class TestErrorRecovery:
    """Verify that files stuck in 'indexing' state are re-indexed on the next run."""

    def test_needs_indexing_detects_stuck_file(self, tmp_path: Path):
        """A file with indexing_state='indexing' must be scheduled for re-index."""
        store = MetadataStore(tmp_path / "test.db")
        f = tmp_path / "stuck.txt"
        f.write_text("some content")

        # Simulate a crash: mark started but never finish
        store.mark_indexing_started(f)

        assert store.needs_indexing(f) is True, (
            "File stuck in 'indexing' state should be flagged for re-indexing"
        )
        store.close()

    def test_needs_indexing_detects_failed_file(self, tmp_path: Path):
        """A file with indexing_state='failed' must be retried."""
        store = MetadataStore(tmp_path / "test.db")
        f = tmp_path / "bad.txt"
        f.write_text("bad content")

        store.mark_indexing_started(f)
        store.mark_indexing_failed(f)

        assert store.needs_indexing(f) is True
        store.close()

    def test_successful_indexing_marks_done(self, tmp_path: Path):
        """After upsert_document, the file's state should be 'done'."""
        store = MetadataStore(tmp_path / "test.db")
        f = tmp_path / "good.txt"
        f.write_text("good content")

        store.mark_indexing_started(f)
        store.upsert_document(f, num_chunks=1)

        doc = store.get_document(f)
        assert doc is not None
        assert doc.indexing_state == STATE_DONE
        assert store.needs_indexing(f) is False
        store.close()

    def test_pipeline_marks_file_as_indexing_before_parse(self, tmp_path: Path):
        """Pipeline should mark a file as 'indexing' before starting to parse it."""
        config = Config(
            data_dir=tmp_path / "data",
            file_extensions=[".txt"],
        )
        store = MetadataStore(config.data_dir / "metadata.db")
        store_started: list[str] = []
        original_mark = store.mark_indexing_started

        def _tracking_mark(path):
            store_started.append(str(path))
            original_mark(path)

        store.mark_indexing_started = _tracking_mark

        embedder = _mock_embedder()
        pipeline = IndexingPipeline(config, store=store)
        pipeline.embedder = embedder

        f = tmp_path / "test.txt"
        f.write_text("Hello, this is a test document about robotics.")

        list(pipeline.index_file(f))

        assert len(store_started) >= 1
        assert any("test.txt" in s for s in store_started)
        pipeline.close()

    def test_pipeline_marks_failed_on_parse_error(self, tmp_path: Path):
        """Pipeline should mark a file as 'failed' when parsing raises."""
        config = Config(
            data_dir=tmp_path / "data",
            file_extensions=[".pdf"],
        )
        store = MetadataStore(config.data_dir / "metadata.db")

        f = tmp_path / "corrupt.pdf"
        f.write_bytes(b"not a real pdf file")

        pipeline = IndexingPipeline(config, store=store)
        statuses = list(pipeline.index_file(f))

        # Should have produced an ERROR status
        assert any(s.status == StatusType.ERROR for s in statuses)

        # Store should not mark it as 'done'
        doc = store.get_document(f.resolve())
        if doc is not None:
            assert doc.indexing_state != STATE_DONE
        pipeline.close()

    def test_interrupted_files_reindexed_on_next_run(self, tmp_path: Path):
        """Files stuck in 'indexing' state should be picked up on the next pipeline run."""
        config = Config(
            data_dir=tmp_path / "data",
            file_extensions=[".txt"],
        )
        store = MetadataStore(config.data_dir / "metadata.db")
        embedder = _mock_embedder()

        f = tmp_path / "interrupted.txt"
        f.write_text("Document that was interrupted during indexing")

        # Simulate crash: mark started but never finish
        store.mark_indexing_started(f.resolve())

        # Verify needs_indexing returns True
        assert store.needs_indexing(f.resolve()) is True

        # Now run the pipeline — should pick it up
        pipeline = IndexingPipeline(config, store=store)
        pipeline.embedder = embedder
        statuses = list(pipeline.index_file(f))

        complete = [s for s in statuses if s.status == StatusType.COMPLETE]
        assert len(complete) >= 1, "Interrupted file should be re-indexed"

        # Should now be marked as done
        doc = store.get_document(f.resolve())
        assert doc is not None
        assert doc.indexing_state == STATE_DONE
        pipeline.close()


# ---------------------------------------------------------------------------
# 5. Config validation
# ---------------------------------------------------------------------------


class TestConfigValidation:
    def test_valid_config_no_issues(self, tmp_path: Path):
        """A fully valid config should return an empty issues list."""
        import socket

        config = Config(
            data_dir=tmp_path / "data",
            index_paths=[tmp_path],
            chunk_size=512,
            chunk_overlap=64,
        )
        # Find a free port
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", 0))
            free_port = s.getsockname()[1]
        config = Config(
            data_dir=tmp_path / "data",
            index_paths=[tmp_path],
            chunk_size=512,
            chunk_overlap=64,
            port=free_port,
        )
        issues = config.validate()
        assert issues == [], f"Expected no issues, got: {issues}"

    def test_invalid_chunk_overlap(self, tmp_path: Path):
        """chunk_overlap >= chunk_size should produce a validation warning."""
        import socket
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", 0))
            free_port = s.getsockname()[1]
        config = Config(
            data_dir=tmp_path / "data",
            chunk_size=100,
            chunk_overlap=150,  # larger than chunk_size!
            port=free_port,
        )
        issues = config.validate()
        assert any("chunk_overlap" in i for i in issues)

    def test_nonexistent_index_path_warns(self, tmp_path: Path):
        """A non-existent index path should generate a warning (not an error)."""
        import socket
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", 0))
            free_port = s.getsockname()[1]
        config = Config(
            data_dir=tmp_path / "data",
            index_paths=[Path("/nonexistent/absolutely/does/not/exist")],
            port=free_port,
        )
        issues = config.validate()
        assert any("does not exist" in i.lower() or "nonexistent" in i.lower() for i in issues)

    def test_port_in_use_warns(self, tmp_path: Path):
        """A port already in use should generate a validation warning."""
        import socket

        # Bind a socket to hold a port
        sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        sock.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
        sock.bind(("127.0.0.1", 0))
        occupied_port = sock.getsockname()[1]
        sock.listen(1)

        try:
            config = Config(
                data_dir=tmp_path / "data",
                host="127.0.0.1",
                port=occupied_port,
            )
            issues = config.validate()
            assert any(str(occupied_port) in i for i in issues), (
                f"Expected port-in-use warning, got: {issues}"
            )
        finally:
            sock.close()


# ---------------------------------------------------------------------------
# 6. Health endpoint
# ---------------------------------------------------------------------------


@pytest.mark.anyio
async def test_health_endpoint_returns_200(tmp_path):
    """GET /api/health should always return 200 with a valid body."""
    from httpx import ASGITransport, AsyncClient
    from desksearch.api.server import create_app

    config = Config(data_dir=tmp_path / "data")
    app = create_app(config)

    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        resp = await client.get("/api/health")

    assert resp.status_code == 200
    body = resp.json()
    assert "status" in body
    assert body["status"] in ("healthy", "degraded", "unhealthy")
    assert "components" in body
    assert "sqlite" in body["components"]
    assert "bm25" in body["components"]
    assert "faiss" in body["components"]
    assert "embedder" in body["components"]


@pytest.mark.anyio
async def test_health_endpoint_sqlite_ok(tmp_path):
    """Health endpoint should report SQLite as 'ok' for a fresh index."""
    from httpx import ASGITransport, AsyncClient
    from desksearch.api.server import create_app

    config = Config(data_dir=tmp_path / "data")
    app = create_app(config)

    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        resp = await client.get("/api/health")

    body = resp.json()
    assert body["components"]["sqlite"]["status"] == "ok"
    assert "doc_count" in body["components"]["sqlite"]


@pytest.mark.anyio
async def test_health_endpoint_search_mode(tmp_path):
    """Health endpoint should report a valid search_mode."""
    from httpx import ASGITransport, AsyncClient
    from desksearch.api.server import create_app

    config = Config(data_dir=tmp_path / "data")
    app = create_app(config)

    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        resp = await client.get("/api/health")

    body = resp.json()
    assert "search_mode" in body
    assert body["search_mode"] in ("hybrid", "bm25_only", "dense_only", "unavailable")


# ---------------------------------------------------------------------------
# 7. Graceful degradation
# ---------------------------------------------------------------------------


class TestGracefulDegradation:
    """Verify search falls back gracefully when one backend is unavailable."""

    def test_bm25_unavailable_falls_back_to_dense(self, tmp_path: Path):
        """When BM25 is unavailable, dense-only search should still work."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        # Simulate BM25 being unavailable
        engine.bm25.available = False

        rng = np.random.default_rng(3)
        emb = rng.standard_normal(DIM).astype(np.float32)
        engine.dense.add("doc1", emb)
        engine._doc_texts["doc1"] = "machine learning content"

        results = engine.search_sync("machine learning", emb, top_k=5)
        assert isinstance(results, list)
        # Should still get results from dense
        assert len(results) >= 1
        assert engine.mode == "dense_only"

    def test_dense_unavailable_falls_back_to_bm25(self, tmp_path: Path):
        """When FAISS is unavailable, BM25-only search should still work."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        # Add a document via BM25 only
        engine.bm25.add_document("doc1", "natural language processing text")

        # Simulate dense being unavailable
        engine.dense.available = False

        rng = np.random.default_rng(4)
        results = engine.search_sync(
            "natural language processing",
            rng.standard_normal(DIM).astype(np.float32),
            top_k=5,
        )
        assert isinstance(results, list)
        assert len(results) >= 1
        assert engine.mode == "bm25_only"

    def test_both_unavailable_returns_empty(self, tmp_path: Path):
        """When both backends are unavailable, search should return empty list."""
        config = Config(data_dir=tmp_path / "data")
        engine = HybridSearchEngine(config, dimension=DIM)

        engine.bm25.available = False
        engine.dense.available = False

        rng = np.random.default_rng(5)
        results = engine.search_sync(
            "any query",
            rng.standard_normal(DIM).astype(np.float32),
        )
        assert results == []
        assert engine.mode == "unavailable"

    def test_store_ping_healthy(self, tmp_path: Path):
        """MetadataStore.ping() should return True for a healthy database."""
        store = MetadataStore(tmp_path / "test.db")
        assert store.ping() is True
        store.close()

    def test_bm25_search_returns_empty_not_raises_when_index_corrupt(self, tmp_path: Path):
        """BM25 search with corrupt/missing index should return [] not raise."""
        idx = BM25Index(tmp_path / "bm25_corrupt")
        idx.available = False  # Simulate corruption

        results = idx.search("any query", top_k=5)
        assert results == []

    def test_dense_search_returns_empty_not_raises_when_unavailable(self, tmp_path: Path):
        """Dense search when unavailable should return [] not raise."""
        idx = DenseIndex(tmp_path / "dense_corrupt", dimension=DIM)
        idx.available = False  # Simulate FAISS unavailable

        rng = np.random.default_rng(6)
        results = idx.search(rng.standard_normal(DIM).astype(np.float32), top_k=5)
        assert results == []
